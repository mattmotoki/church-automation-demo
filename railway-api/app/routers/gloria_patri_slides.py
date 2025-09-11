"""
Gloria Patri slide generation endpoints
"""
import os
import tempfile
import base64
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from app.core.schemas import GloriaPatriRequest
from pathlib import Path

router = APIRouter()


def create_gloria_patri_slide(output_file, background_image=None):
    """Create a gloria patri slide"""
    # Get the template path
    template_path = Path(__file__).parent.parent / "templates" / "gloria-patri.pptx"
    if not template_path.exists():
        raise ValueError(f"Template not found: {template_path}")
    
    
    # Load the template presentation
    prs = Presentation(template_path)
    
    # Add background image if provided
    if background_image and len(prs.slides) > 0:
        slide = prs.slides[0]
        try:
            # Decode base64 image
            image_data = base64.b64decode(background_image.split(',')[1] if ',' in background_image else background_image)
            
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
                tmp_img.write(image_data)
                tmp_img_path = tmp_img.name
            
            # Get slide dimensions
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Add the background image covering the entire slide
            pic = slide.shapes.add_picture(tmp_img_path, 0, 0, slide_width, slide_height)
            
            # Move the picture to the back (z-order)
            # First remove it from its current position (at the end)
            slide.shapes._spTree.remove(pic._element)
            # Then insert it at position 2 (after nvGrpSpPr at 0 and grpSpPr at 1)
            slide.shapes._spTree.insert(2, pic._element)
            
            # Clean up temp file
            os.unlink(tmp_img_path)
            
        except Exception as e:
            print(f"Error processing background image: {str(e)}")
    
    # Save the presentation
    prs.save(output_file)
    return output_file


@router.post("/api/generate-gloria-patri-slide")
async def generate_gloria_patri_slide(request: GloriaPatriRequest):
    """Generate a Gloria Patri slide with a background image"""
    try:
        # Get the template path
        template_path = Path(__file__).parent.parent / "templates" / "gloria-patri.pptx"
        if not template_path.exists():
            raise HTTPException(status_code=500, detail=f"Template not found: {template_path}")
        
        
        # Load the template presentation
        prs = Presentation(template_path)
        
        # Get the first slide
        if len(prs.slides) == 0:
            raise HTTPException(
                status_code=500,
                detail="Template has no slides"
            )
        
        slide = prs.slides[0]
        
        # Add the background image as a new full-slide shape at the back
        if request.background_image:
            try:
                # Decode base64 image
                image_data = base64.b64decode(request.background_image.split(',')[1] if ',' in request.background_image else request.background_image)
                
                # Save to temporary file
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
                    tmp_img.write(image_data)
                    tmp_img_path = tmp_img.name
                
                # Get slide dimensions
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # Add the background image covering the entire slide
                pic = slide.shapes.add_picture(tmp_img_path, 0, 0, slide_width, slide_height)
                
                # Move the picture to the back (z-order)
                # First remove it from its current position (at the end)
                slide.shapes._spTree.remove(pic._element)
                # Then insert it at position 2 (after nvGrpSpPr at 0 and grpSpPr at 1)
                slide.shapes._spTree.insert(2, pic._element)
                
                # Clean up temp file
                os.unlink(tmp_img_path)
                
            except Exception as e:
                print(f"Error processing background image: {str(e)}")
                # Continue without adding the background
        
        # Save to a temporary file
        with tempfile.NamedTemporaryFile(
            suffix='.pptx',
            delete=False,
            prefix='gloria_patri_slide_'
        ) as tmp_file:
            prs.save(tmp_file.name)
            tmp_path = tmp_file.name
        
        # Return the file
        return FileResponse(
            tmp_path,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            filename='gloria_patri_slide.pptx',
            headers={
                'Content-Disposition': 'attachment; filename="gloria_patri_slide.pptx"'
            }
        )
        
    except Exception as e:
        print(f"Error generating Gloria Patri slide: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Error generating Gloria Patri slide: {str(e)}"
        )