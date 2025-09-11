"""
Welcome slide generation endpoints
"""
import os
import tempfile
from pathlib import Path
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Pt
from app.core.schemas import WelcomeSlideRequest

router = APIRouter()


@router.post("/api/generate-welcome-slide")
async def generate_welcome_slide(request: WelcomeSlideRequest):
    """Generate a welcome slide with the lead pastor's name"""
    try:
        # Get the template path (relative to this file)
        template_path = Path(__file__).parent.parent / "templates" / "welcome.pptx"
        
        if not template_path.exists():
            raise HTTPException(status_code=500, detail=f"Template not found: {template_path}")
        
        # Load the template presentation
        prs = Presentation(template_path)
        
        # Get the first slide
        if len(prs.slides) == 0:
            raise HTTPException(
                status_code=500,
                detail="Template has no slides"
            )
        
        slide = prs.slides[0]
        
        # Find the text shape (first shape is the text box)
        if len(slide.shapes) < 1:
            raise HTTPException(
                status_code=500,
                detail="Template slide has no shapes"
            )
        
        text_shape = slide.shapes[0]
        
        # Replace the {lead_pastor} placeholder
        if hasattr(text_shape, 'text_frame'):
            # Find and replace the placeholder while preserving formatting
            for paragraph in text_shape.text_frame.paragraphs:
                # Get the full text from all runs combined
                full_text = ''.join(run.text for run in paragraph.runs) if paragraph.runs else paragraph.text
                
                # Check if this paragraph has the placeholder
                if full_text and '{lead_pastor}' in full_text:
                    # Replace the placeholder in the full text
                    new_text = full_text.replace('{lead_pastor}', request.lead_pastor)
                    
                    # If there are multiple runs (text split across formatting), handle specially
                    if len(paragraph.runs) >= 3:
                        # Assume the placeholder is split across 3 runs: "{", "lead_pastor", "}"
                        # We'll modify them in place to preserve formatting
                        found_placeholder = False
                        for i in range(len(paragraph.runs) - 2):
                            if (paragraph.runs[i].text == '{' and 
                                paragraph.runs[i+1].text == 'lead_pastor' and 
                                paragraph.runs[i+2].text == '}'):
                                # Replace the middle run with the lead pastor name
                                paragraph.runs[i+1].text = request.lead_pastor
                                paragraph.runs[i].text = ""
                                paragraph.runs[i+2].text = ""
                                found_placeholder = True
                                break
                        
                        if not found_placeholder:
                            # If not split in expected way, replace in first run
                            if paragraph.runs:
                                paragraph.runs[0].text = new_text
                                # Clear other runs
                                for i in range(1, len(paragraph.runs)):
                                    paragraph.runs[i].text = ""
                    elif paragraph.runs:
                        # Single or few runs - just replace in the first run
                        paragraph.runs[0].text = new_text
                        # Clear any other runs
                        for i in range(1, len(paragraph.runs)):
                            paragraph.runs[i].text = ""
                    else:
                        # No runs, set paragraph text directly
                        paragraph.text = new_text
                    break
        
        # Save to a temporary file
        with tempfile.NamedTemporaryFile(
            suffix='.pptx',
            delete=False,
            prefix='welcome_slide_'
        ) as tmp_file:
            prs.save(tmp_file.name)
            tmp_path = tmp_file.name
        
        # Return the file
        return FileResponse(
            tmp_path,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            filename='welcome_slide.pptx',
            headers={
                'Content-Disposition': 'attachment; filename="welcome_slide.pptx"'
            }
        )
        
    except Exception as e:
        print(f"Error generating welcome slide: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Error generating welcome slide: {str(e)}"
        )