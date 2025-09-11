"""
Postlude slides router for generating PowerPoint presentations
"""
import os
import shutil
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from typing import Dict, Any
from app.core.files import create_temp_file
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from .slides.utils import (
    process_background_image,
    cleanup_temp_file
)

router = APIRouter()


def create_postlude_slides(output_file, pianist_name="Jason Eom", background_image=None):
    """Create postlude slides from template"""
    # Path to the template file
    template_path = Path(__file__).parent.parent / "templates" / "postlude.pptx"
    if not template_path.exists():
        raise ValueError(f"Template not found: {template_path}")
    
    
    # Load template presentation
    prs = Presentation(template_path)
    
    # Update pianist name in slide 2 (index 1)
    if len(prs.slides) > 1:
        slide = prs.slides[1]
        # Find text box with pianist placeholder
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    full_text = ''.join(run.text for run in paragraph.runs)
                    if '{pianist}' in full_text:
                        if len(paragraph.runs) >= 3:
                            paragraph.runs[1].text = pianist_name
                            paragraph.runs[0].text = ""
                            paragraph.runs[2].text = ""
                        elif len(paragraph.runs) == 1:
                            paragraph.runs[0].text = paragraph.runs[0].text.replace('{pianist}', pianist_name)
                        break
    
    # Add background image if provided (base64)
    if background_image:
        try:
            # Decode base64 image
            image_data = base64.b64decode(background_image.split(',')[1] if ',' in background_image else background_image)
            
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
                tmp_img.write(image_data)
                tmp_img_path = tmp_img.name
            
            # Add background to each slide
            for slide in prs.slides:
                # Get slide dimensions
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # Add the background image covering the entire slide
                pic = slide.shapes.add_picture(tmp_img_path, 0, 0, slide_width, slide_height)
                
                # Move the picture to the back (z-order)
                # First remove it from its current position (at the end)
                slide.shapes._spTree.remove(pic._element)
                # Then insert it at position 2 (after nvGrpSpPr at 0 and grpSpPr at 1)
                slide.shapes._spTree.insert(2, pic._element)
            
            # Clean up temp file
            os.unlink(tmp_img_path)
            
        except Exception as e:
            print(f"Error processing background image: {str(e)}")
    
    # Save the presentation
    prs.save(output_file)
    return output_file


@router.post("/generate-postlude-slides")
async def generate_postlude_slides(data: Dict[str, Any]):
    """Generate postlude PowerPoint slides from template with optional background and pianist name"""
    try:
        # Get pianist name from request, default to "Jason Eom"
        pianist_name = data.get('pianist', 'Jason Eom').strip()
        if not pianist_name:
            pianist_name = 'Jason Eom'
        
        # Path to the template file
        template_path = Path(__file__).parent.parent / "templates" / "postlude.pptx"
        if not template_path.exists():
            raise HTTPException(status_code=500, detail=f"Template not found: {template_path}")
        
        # Create temporary output file
        output_path = create_temp_file(suffix='.pptx')
        
        # Copy the template to the output file
        shutil.copy2(template_path, output_path)
        
        # Open the presentation to modify it
        prs = Presentation(output_path)
        
        # Process background image if provided
        background_image = data.get('background_image')
        if background_image:
            # Process background image - expects base64
            background_image_path = process_background_image(background_image)
            
            # Apply background to all slides by inserting image at the back
            for slide in prs.slides:
                # Add background image as the first shape (behind all other content)
                left = 0
                top = 0
                pic = slide.shapes.add_picture(
                    background_image_path, 
                    left, 
                    top,
                    prs.slide_width,
                    prs.slide_height
                )
                # Move the picture to the back by rearranging the shape tree
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)  # Insert after nvGrpSpPr and grpSpPr
            
            # Clean up temporary background image file
            cleanup_temp_file(background_image_path)
        
        # Update pianist name in the slides while preserving formatting
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # Check each paragraph for the pianist name
                    for paragraph in shape.text_frame.paragraphs:
                        # Check if this paragraph contains "Jason Eom"
                        full_text = ''.join(run.text for run in paragraph.runs)
                        if 'Jason Eom' in full_text:
                            # Replace while preserving formatting by modifying runs in place
                            for run in paragraph.runs:
                                if 'Jason Eom' in run.text:
                                    # Replace text in this run while keeping all formatting
                                    run.text = run.text.replace('Jason Eom', pianist_name)
        
        # Save the modified presentation
        prs.save(output_path)
        
        # Return the file
        return FileResponse(
            path=output_path,
            filename="postlude_slides.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))