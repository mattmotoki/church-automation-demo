"""
Combined slides router for generating a single PowerPoint presentation with all service elements
"""
import os
import zipfile
import shutil
import xml.etree.ElementTree as ET
import re
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from typing import Dict, Any, List, Optional
from app.core.files import create_temp_file
from pathlib import Path
from .slides.utils import cleanup_temp_file

router = APIRouter()


def merge_pptx_simple(pptx_files: List[str], slides_to_take: List[int] = None) -> str:
    """Merge PPTX files by extracting and recombining their contents"""
    if not pptx_files:
        raise ValueError("No presentations to merge")
    
    # Create a unique temporary directory for extraction
    import tempfile
    temp_dir = tempfile.mkdtemp(suffix='_merge_dir')  # Creates unique temp dir
    
    try:
        # Extract the first presentation as base, but we'll clear its slides
        print(f"Using {os.path.basename(pptx_files[0])} as base template")
        base_dir = os.path.join(temp_dir, 'base')
        with zipfile.ZipFile(pptx_files[0], 'r') as zip_ref:
            zip_ref.extractall(base_dir)
        
        # Parse presentation.xml
        pres_xml_path = os.path.join(base_dir, 'ppt', 'presentation.xml')
        tree = ET.parse(pres_xml_path)
        root = tree.getroot()
        
        # Define namespaces
        ns = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        }
        
        # Register namespaces to preserve them in output
        for prefix, uri in ns.items():
            ET.register_namespace(prefix, uri)
        
        # Find slide list and CLEAR it - we'll rebuild from scratch
        sld_id_list = root.find('.//p:sldIdLst', ns)
        if sld_id_list is None:
            sld_id_list = ET.SubElement(root, '{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
        else:
            # Clear all existing slides from the list
            for child in list(sld_id_list):
                sld_id_list.remove(child)
            print(f"  Cleared existing slide list")
        
        # Start fresh with IDs
        max_slide_id = 256
        max_rel_id = 10  # Start higher to avoid conflicts with theme/layout relationships
        
        # Clear the slides directory
        base_slides_dir = os.path.join(base_dir, 'ppt', 'slides')
        if os.path.exists(base_slides_dir):
            # Remove all existing slide files
            for slide_file in os.listdir(base_slides_dir):
                if slide_file.startswith('slide') and slide_file.endswith('.xml'):
                    os.remove(os.path.join(base_slides_dir, slide_file))
            print(f"  Cleared existing slides from base")
        
        # Also clear slide relationships
        base_slides_rels_dir = os.path.join(base_slides_dir, '_rels')
        if os.path.exists(base_slides_rels_dir):
            import shutil
            shutil.rmtree(base_slides_rels_dir)
        os.makedirs(base_slides_rels_dir, exist_ok=True)
        
        slide_count = 0
        print(f"  Starting with clean slate (0 slides)")
        
        # Parse presentation.xml.rels and remove existing slide relationships
        pres_rels_path = os.path.join(base_dir, 'ppt', '_rels', 'presentation.xml.rels')
        rels_tree = ET.parse(pres_rels_path)
        rels_root = rels_tree.getroot()
        
        # Remove any existing slide relationships
        for rel in list(rels_root):
            if 'slides/slide' in rel.get('Target', ''):
                rels_root.remove(rel)
        print(f"  Cleared slide relationships")
        
        # Process ALL presentations (including the first one)
        for pptx_idx, pptx_file in enumerate(pptx_files):
            print(f"\n--- Processing presentation {pptx_idx}: {os.path.basename(pptx_file)} ---")
            source_dir = os.path.join(temp_dir, f'source_{pptx_idx}')
            
            # Verify file exists and is valid
            if not os.path.exists(pptx_file):
                print(f"  WARNING: File does not exist: {pptx_file}")
                continue
                
            with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
                zip_ref.extractall(source_dir)
            
            # Track relationship files that need media path updates
            rels_files_to_update = []
            media_mapping = {}  # Map old names to new names
            
            # Get source slides
            source_slides_dir = os.path.join(source_dir, 'ppt', 'slides')
            if os.path.exists(source_slides_dir):
                source_slides = sorted([f for f in os.listdir(source_slides_dir) 
                                       if f.startswith('slide') and f.endswith('.xml')],
                                     key=lambda x: int(re.findall(r'\d+', x)[0]))
                
                print(f"  Found {len(source_slides)} slides to add")
                if len(source_slides) == 0:
                    print(f"  WARNING: No slides found in {os.path.basename(pptx_file)}")
                    continue
                
                for slide_file in source_slides:
                    source_slide_num = int(re.findall(r'\d+', slide_file)[0])
                    slide_count += 1
                    max_rel_id += 1
                    max_slide_id += 1
                    
                    # Copy slide file
                    src_slide = os.path.join(source_slides_dir, slide_file)
                    dst_slide = os.path.join(base_slides_dir, f'slide{slide_count}.xml')
                    shutil.copy2(src_slide, dst_slide)
                    
                    # Copy slide relationships if exists
                    src_rels_dir = os.path.join(source_slides_dir, '_rels')
                    if os.path.exists(src_rels_dir):
                        src_rels_file = os.path.join(src_rels_dir, f'slide{source_slide_num}.xml.rels')
                        if os.path.exists(src_rels_file):
                            dst_rels_dir = os.path.join(base_slides_dir, '_rels')
                            os.makedirs(dst_rels_dir, exist_ok=True)
                            dst_rels_file = os.path.join(dst_rels_dir, f'slide{slide_count}.xml.rels')
                            
                            # Read and update the relationships file
                            with open(src_rels_file, 'r', encoding='utf-8') as f:
                                rels_content = f.read()
                            
                            # Track this file for media reference updates
                            rels_files_to_update.append(dst_rels_file)
                            
                            # Write the relationships file (will be updated after media copy)
                            with open(dst_rels_file, 'w', encoding='utf-8') as f:
                                f.write(rels_content)
                    
                    # Add slide to presentation.xml
                    new_sld_id = ET.SubElement(sld_id_list, '{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
                    new_sld_id.set('id', str(max_slide_id))
                    new_sld_id.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', f'rId{max_rel_id}')
                    
                    # Add relationship to presentation.xml.rels
                    new_rel = ET.SubElement(rels_root, 'Relationship')
                    new_rel.set('Id', f'rId{max_rel_id}')
                    new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
                    new_rel.set('Target', f'slides/slide{slide_count}.xml')
            
            else:
                print(f"  WARNING: No slides directory found in {os.path.basename(pptx_file)}")
                # Check what directories exist
                ppt_dir = os.path.join(source_dir, 'ppt')
                if os.path.exists(ppt_dir):
                    print(f"    Contents of ppt/: {os.listdir(ppt_dir)}")
                else:
                    print(f"    ERROR: No ppt directory found")
                continue
            
            # Copy media files with unique names to avoid conflicts
            source_media_dir = os.path.join(source_dir, 'ppt', 'media')
            if os.path.exists(source_media_dir):
                base_media_dir = os.path.join(base_dir, 'ppt', 'media')
                os.makedirs(base_media_dir, exist_ok=True)
                
                for media_file in os.listdir(source_media_dir):
                    src = os.path.join(source_media_dir, media_file)
                    # Create unique name by prefixing with presentation index
                    name, ext = os.path.splitext(media_file)
                    new_media_file = f"pres{pptx_idx}_{media_file}"
                    dst = os.path.join(base_media_dir, new_media_file)
                    shutil.copy2(src, dst)
                    media_mapping[media_file] = new_media_file
            
            # Update relationship files with new media paths
            for rels_file in rels_files_to_update:
                if os.path.exists(rels_file):
                    with open(rels_file, 'r', encoding='utf-8') as f:
                        rels_content = f.read()
                    
                    # Update media references
                    for old_name, new_name in media_mapping.items():
                        rels_content = rels_content.replace(f'Target="../media/{old_name}"', 
                                                          f'Target="../media/{new_name}"')
                    
                    with open(rels_file, 'w', encoding='utf-8') as f:
                        f.write(rels_content)
            
            # Copy slide layouts if missing
            source_layouts_dir = os.path.join(source_dir, 'ppt', 'slideLayouts')
            if os.path.exists(source_layouts_dir):
                base_layouts_dir = os.path.join(base_dir, 'ppt', 'slideLayouts')
                if not os.path.exists(base_layouts_dir):
                    shutil.copytree(source_layouts_dir, base_layouts_dir)
        
        # Write updated presentation.xml
        tree.write(pres_xml_path, encoding='utf-8', xml_declaration=True)
        
        # Write updated presentation.xml.rels
        rels_tree.write(pres_rels_path, encoding='utf-8', xml_declaration=True)
        
        # Fix Content Types if needed
        content_types_path = os.path.join(base_dir, '[Content_Types].xml')
        if os.path.exists(content_types_path):
            ct_tree = ET.parse(content_types_path)
            ct_root = ct_tree.getroot()
            
            # Ensure slide content type is registered
            ct_ns = {'': 'http://schemas.openxmlformats.org/package/2006/content-types'}
            
            # Check if slide override exists for all slides
            for i in range(1, slide_count + 1):
                slide_path = f'/ppt/slides/slide{i}.xml'
                override_exists = False
                
                for override in ct_root.findall('Override', ct_ns):
                    if override.get('PartName') == slide_path:
                        override_exists = True
                        break
                
                if not override_exists:
                    new_override = ET.SubElement(ct_root, 'Override')
                    new_override.set('PartName', slide_path)
                    new_override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
            
            ct_tree.write(content_types_path, encoding='utf-8', xml_declaration=True)
        
        # Create output PPTX
        print(f"\n=== MERGE COMPLETE ===")
        print(f"Total slides in merged presentation: {slide_count}")
        
        # List all slides in the final presentation
        final_slides = sorted([f for f in os.listdir(base_slides_dir) 
                              if f.startswith('slide') and f.endswith('.xml')],
                             key=lambda x: int(re.findall(r'\d+', x)[0]))
        print(f"Final slide files: {final_slides[:10]}...") # Show first 10
        
        output_path = create_temp_file(suffix='.pptx')
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(base_dir):
                for file in files:
                    if file != '.DS_Store':  # Skip macOS files
                        file_path = os.path.join(root, file)
                        arc_path = os.path.relpath(file_path, base_dir)
                        zipf.write(file_path, arc_path)
        
        return output_path
        
    finally:
        # Clean up temp directory
        shutil.rmtree(temp_dir, ignore_errors=True)


@router.post("/generate-combined-slides")
async def generate_combined_slides(data: Dict[str, Any]):
    """Generate a single PowerPoint presentation with all service elements"""
    temp_files = []
    
    try:
        slides_added = []
        
        # Import needed for template paths
        import shutil
        from pptx import Presentation
        from pptx.util import Pt
        
        # Generate all slide sets in order
        
        # 1. Shell Blowing
        # Copy template directly
        try:
            shell_template = Path(__file__).parent.parent / "templates" / "shell-blowing.pptx"
            if not shell_template.exists():
                raise FileNotFoundError(f"Template not found: {shell_template}")
            shell_output = create_temp_file(suffix='.pptx')
            import shutil
            shutil.copy2(shell_template, shell_output)
            temp_files.append(shell_output)
            slides_added.append("Shell Blowing")
        except FileNotFoundError as e:
            print(f"Warning: Shell blowing template not found: {e}")
        
        # 2. Welcome Slide
        # Create a welcome slide function to properly generate it
        def create_welcome_slide(output_file, lead_pastor="Pastor"):
            """Create a welcome slide from template"""
            try:
                welcome_template = Path(__file__).parent.parent / "templates" / "welcome.pptx"
                if not welcome_template.exists():
                    raise FileNotFoundError(f"Template not found: {welcome_template}")
            except FileNotFoundError as e:
                raise ValueError(str(e))
            
            from pptx import Presentation
            prs = Presentation(welcome_template)
            
            # Get only the first slide (template might have multiple)
            if not prs.slides:
                raise ValueError("Welcome template has no slides")
            
            # Create a new presentation with just the first slide
            new_prs = Presentation()
            
            # Copy the first slide's layout
            slide_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None
            
            # For now, let's just use the template as-is but ensure we only have one slide
            # We'll modify the text in the first slide
            slide = prs.slides[0]
            if slide.shapes:
                text_shape = slide.shapes[0]
                if hasattr(text_shape, 'text_frame'):
                    for paragraph in text_shape.text_frame.paragraphs:
                        # Get the full text from all runs combined
                        full_text = ''.join(run.text for run in paragraph.runs) if paragraph.runs else paragraph.text
                        
                        # Check if this paragraph has the placeholder
                        if full_text and '{lead_pastor}' in full_text:
                            # Replace the placeholder in the full text
                            new_text = full_text.replace('{lead_pastor}', lead_pastor)
                            
                            # Handle split runs
                            if len(paragraph.runs) >= 3:
                                found_placeholder = False
                                for i in range(len(paragraph.runs) - 2):
                                    if (paragraph.runs[i].text == '{' and 
                                        paragraph.runs[i+1].text == 'lead_pastor' and 
                                        paragraph.runs[i+2].text == '}'):
                                        paragraph.runs[i+1].text = lead_pastor
                                        paragraph.runs[i].text = ""
                                        paragraph.runs[i+2].text = ""
                                        found_placeholder = True
                                        break
                                
                                if not found_placeholder and paragraph.runs:
                                    paragraph.runs[0].text = new_text
                                    for i in range(1, len(paragraph.runs)):
                                        paragraph.runs[i].text = ""
                            elif paragraph.runs:
                                paragraph.runs[0].text = new_text
                                for i in range(1, len(paragraph.runs)):
                                    paragraph.runs[i].text = ""
                            else:
                                paragraph.text = new_text
                            break
            
            # Save only the first slide (create new presentation to avoid issues)
            # This is a workaround - save the whole thing but only use first slide in merge
            print(f"Welcome template has {len(prs.slides)} slides")
            prs.save(output_file)
            
            # Verify the saved file
            verify_prs = Presentation(output_file)
            print(f"Saved welcome file has {len(verify_prs.slides)} slides")
            
            return output_file
        
        # Generate welcome slide
        welcome_output = create_temp_file(suffix='.pptx')
        try:
            create_welcome_slide(welcome_output, data.get('lead_pastor', 'Pastor'))
            temp_files.append(welcome_output)
            slides_added.append("Welcome")
        except Exception as e:
            print(f"Error generating welcome slide: {e}")
        
        # 3. Call to Worship
        if data.get('call_to_worship_pairs'):
            print(f"\nGenerating Call to Worship slides...")
            print(f"  Number of pairs: {len(data['call_to_worship_pairs'])}")
            ctw_output = create_temp_file(suffix='.pptx')
            from .call_to_worship_slides import create_call_to_worship_slides_from_dict
            try:
                create_call_to_worship_slides_from_dict(
                    data['call_to_worship_pairs'],
                    ctw_output,
                    data.get('base_background')
                )
                
                # Verify the created file
                from pptx import Presentation
                verify_prs = Presentation(ctw_output)
                print(f"  Created Call to Worship presentation with {len(verify_prs.slides)} slides")
                
                temp_files.append(ctw_output)
                slides_added.append("Call to Worship")
            except Exception as e:
                print(f"  ERROR generating Call to Worship slides: {e}")
                import traceback
                traceback.print_exc()
        else:
            print(f"\nNo Call to Worship pairs provided, skipping...")
        
        # 4. Hymns (up to 3)
        hymns = data.get('hymns', [])
        for i, hymn_data in enumerate(hymns[:3]):
            if hymn_data and hymn_data.get('hymn'):
                hymn_output = create_temp_file(suffix='.pptx')
                from .hymn_slides import create_hymn_slides
                bg_key = f'hymn{i + 1}_background'
                background = data.get(bg_key)
                # Transform hymn data to match expected format
                hymn_info = hymn_data['hymn']
                hymn_info['hymn_number'] = hymn_info.get('number', '')
                if 'lyrics' not in hymn_info:
                    hymn_info['lyrics'] = []
                create_hymn_slides(
                    hymn_info,
                    hymn_output,
                    background
                )
                temp_files.append(hymn_output)
                slides_added.append(f"Hymn {i+1}")
        
        # 5. Message for All Generations
        try:
            message_template = Path(__file__).parent.parent / "templates" / "message-for-all-generations.pptx"
            if not message_template.exists():
                raise FileNotFoundError(f"Template not found: {message_template}")
            from pptx import Presentation
            from pptx.util import Pt
            prs = Presentation(message_template)
            if prs.slides and len(prs.slides[0].shapes) >= 2:
                text_shape = prs.slides[0].shapes[1]
                if hasattr(text_shape, 'text_frame'):
                    # Find the paragraph with "{lead_pastor}" placeholder
                    for paragraph in text_shape.text_frame.paragraphs:
                        # Check if this paragraph contains the placeholder
                        full_text = ''.join(run.text for run in paragraph.runs)
                        if '{lead_pastor}' in full_text:
                            # The text is split across 3 runs: "{", "lead_pastor", "}"
                            # We'll modify them in place to preserve all XML properties including glow effect
                            
                            if len(paragraph.runs) >= 3:
                                # Find the runs that contain the placeholder
                                for i in range(len(paragraph.runs) - 2):
                                    if (paragraph.runs[i].text == '{' and 
                                        paragraph.runs[i+1].text == 'lead_pastor' and 
                                        paragraph.runs[i+2].text == '}'):
                                        # Replace the middle run with the lead pastor name
                                        paragraph.runs[i+1].text = data.get('lead_pastor', 'Pastor')
                                        paragraph.runs[i].text = ""
                                        paragraph.runs[i+2].text = ""
                                        break
                                else:
                                    # If not found in expected format, try simple replacement
                                    if paragraph.runs:
                                        paragraph.runs[0].text = full_text.replace('{lead_pastor}', data.get('lead_pastor', 'Pastor'))
                                        for i in range(1, len(paragraph.runs)):
                                            paragraph.runs[i].text = ""
                            elif len(paragraph.runs) == 1:
                                # If it's a single run, just replace the text
                                paragraph.runs[0].text = paragraph.runs[0].text.replace('{lead_pastor}', data.get('lead_pastor', 'Pastor'))
                            break
            message_output = create_temp_file(suffix='.pptx')
            prs.save(message_output)
            temp_files.append(message_output)
            slides_added.append("Message for All Generations")
        except FileNotFoundError as e:
            print(f"Warning: Message for All Generations template not found: {e}")
        
        # 6. Prayer of Dedication
        from .prayer_of_dedication_slides import create_prayer_of_dedication_slide
        prayer_output = create_temp_file(suffix='.pptx')
        create_prayer_of_dedication_slide(
            prayer_output,
            data.get('base_background'),
            data.get('lead_pastor', 'Pastor')
        )
        temp_files.append(prayer_output)
        slides_added.append("Prayer of Dedication")
        
        # 7. Gloria Patri  
        from .gloria_patri_slides import create_gloria_patri_slide
        gloria_output = create_temp_file(suffix='.pptx')
        create_gloria_patri_slide(
            gloria_output,
            data.get('base_background')
        )
        temp_files.append(gloria_output)
        slides_added.append("Gloria Patri")
        
        # 8. Scripture
        if data.get('scripture'):
            from .scripture_slides import create_scripture_slides
            scripture_output = create_temp_file(suffix='.pptx')
            create_scripture_slides(
                data['scripture'].get('reference', {}),
                data['scripture'].get('verses', []),
                scripture_output,
                data.get('base_background')
            )
            temp_files.append(scripture_output)
            slides_added.append("Scripture")
        
        # 9. Doxology
        from .doxology_slides import create_doxology_slides
        doxology_output = create_temp_file(suffix='.pptx')
        create_doxology_slides(
            doxology_output,
            data.get('base_background')
        )
        temp_files.append(doxology_output)
        slides_added.append("Doxology")
        
        # 10. Benediction
        from .benediction_slides import create_benediction_slide
        benediction_output = create_temp_file(suffix='.pptx')
        create_benediction_slide(
            benediction_output,
            data.get('base_background'),
            data.get('lead_pastor', 'Pastor')
        )
        temp_files.append(benediction_output)
        slides_added.append("Benediction")
        
        # 11. Postlude
        from .postlude_slides import create_postlude_slides
        postlude_output = create_temp_file(suffix='.pptx')
        create_postlude_slides(
            postlude_output,
            data.get('pianist_name', 'Jason Eom'),
            data.get('base_background')
        )
        temp_files.append(postlude_output)
        slides_added.append("Postlude")
        
        # Combine all presentations using XML manipulation
        if temp_files:
            print(f"\n=== READY TO MERGE ===")
            print(f"Total presentations to merge: {len(temp_files)}")
            print(f"Components: {slides_added}")
            
            # Verify each file before merging
            for i, (temp_file, component) in enumerate(zip(temp_files, slides_added)):
                try:
                    from pptx import Presentation
                    prs = Presentation(temp_file)
                    print(f"  {i}. {component}: {len(prs.slides)} slides - {temp_file}")
                except Exception as e:
                    print(f"  {i}. {component}: ERROR reading file - {e}")
            
            combined_file = merge_pptx_simple(temp_files)
            
            # Clean up temp files
            for temp_file in temp_files:
                try:
                    cleanup_temp_file(temp_file)
                except:
                    pass
            
            # Return the combined file
            return FileResponse(
                path=combined_file,
                filename="complete_service_slides.pptx",
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={
                    "X-Slides-Added": ",".join(slides_added)
                }
            )
        else:
            raise HTTPException(status_code=400, detail="No slides were generated")
        
    except Exception as e:
        # Clean up any temp files on error
        for temp_file in temp_files:
            try:
                cleanup_temp_file(temp_file)
            except:
                pass
        
        raise HTTPException(status_code=500, detail=str(e))