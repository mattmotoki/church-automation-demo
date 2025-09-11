"""
Doxology slides router for generating PowerPoint presentations
"""
import os
import shutil
import base64
import tempfile
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from typing import Dict, Any
from app.core.files import create_temp_file
from pathlib import Path
from pptx import Presentation
from .slides.utils import (
    process_background_image,
    set_slide_background,
    cleanup_temp_file
)

router = APIRouter()


def create_doxology_slides(output_file, background_image=None):
    """Create doxology slides from template"""
    import base64
    import tempfile
    
    # Path to the template file
    template_path = Path(__file__).parent.parent / "templates" / "doxology.pptx"
    if not template_path.exists():
        raise ValueError(f"Template not found: {template_path}")
    
    
    # Load template presentation
    prs = Presentation(template_path)
    
    # Add background image if provided (base64)
    if background_image:
        try:
            # Decode base64 image
            image_data = base64.b64decode(background_image.split(',')[1] if ',' in background_image else background_image)
            
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
                tmp_img.write(image_data)
                tmp_img_path = tmp_img.name
            
            # Add background to each slide
            for slide in prs.slides:
                # Get slide dimensions
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # Add the background image covering the entire slide
                pic = slide.shapes.add_picture(tmp_img_path, 0, 0, slide_width, slide_height)
                
                # Move the picture to the back (z-order)
                # First remove it from its current position (at the end)
                slide.shapes._spTree.remove(pic._element)
                # Then insert it at position 2 (after nvGrpSpPr at 0 and grpSpPr at 1)
                slide.shapes._spTree.insert(2, pic._element)
            
            # Clean up temp file
            os.unlink(tmp_img_path)
            
        except Exception as e:
            print(f"Error processing background image: {str(e)}")
    
    # Save the presentation
    prs.save(output_file)
    return output_file


@router.post("/generate-doxology-slides")
async def generate_doxology_slides(data: Dict[str, Any]):
    """Generate doxology PowerPoint slides from template with optional background"""
    try:
        # Path to the template file
        template_path = Path(__file__).parent.parent / "templates" / "doxology.pptx"
        if not template_path.exists():
            raise HTTPException(status_code=500, detail=f"Template not found: {template_path}")
        
        # Create temporary output file
        output_path = create_temp_file(suffix='.pptx')
        
        # Copy the template to the output file
        shutil.copy2(template_path, output_path)
        
        # Process background image if provided
        background_image = data.get('background_image')
        if background_image:
            # Open the presentation
            prs = Presentation(output_path)
            
            # Process background image - expects base64
            background_image_path = process_background_image(background_image)
            
            # Apply background to all slides by inserting image at the back
            for slide in prs.slides:
                # Add background image as the first shape (behind all other content)
                left = 0
                top = 0
                pic = slide.shapes.add_picture(
                    background_image_path, 
                    left, 
                    top,
                    prs.slide_width,
                    prs.slide_height
                )
                # Move the picture to the back by rearranging the shape tree
                # The picture was just added, so it's the last shape
                # We need to move it to be the first shape (index 0)
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)  # Insert after nvGrpSpPr and grpSpPr
            
            # Save the modified presentation
            prs.save(output_path)
            
            # Clean up temporary background image file
            cleanup_temp_file(background_image_path)
        
        # Return the file
        return FileResponse(
            path=output_path,
            filename="doxology_slides.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))