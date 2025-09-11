"""
Message for All Generations slide generation endpoints
"""
import os
import tempfile
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Pt
from app.core.schemas import MessageForAllGenerationsRequest
from pathlib import Path

router = APIRouter()


@router.post("/api/generate-message-for-all-generations-slide")
async def generate_message_for_all_generations_slide(request: MessageForAllGenerationsRequest):
    """Generate a Message for All Generations slide with the lead pastor's name"""
    try:
        # Get the template path
        template_path = Path(__file__).parent.parent / "templates" / "message-for-all-generations.pptx"
        if not template_path.exists():
            raise HTTPException(status_code=500, detail=f"Template not found: {template_path}")
        
        
        # Load the template presentation
        prs = Presentation(template_path)
        
        # Get the first slide
        if len(prs.slides) == 0:
            raise HTTPException(
                status_code=500,
                detail="Template has no slides"
            )
        
        slide = prs.slides[0]
        
        # Find the text shape (shape at index 1, after the picture)
        if len(slide.shapes) < 2:
            raise HTTPException(
                status_code=500,
                detail="Template slide doesn't have expected shapes"
            )
        
        text_shape = slide.shapes[1]
        
        # Replace the {lead_pastor} placeholder while preserving formatting
        if hasattr(text_shape, 'text_frame'):
            # Find the paragraph with "{lead_pastor}" text
            for paragraph in text_shape.text_frame.paragraphs:
                # Check if this paragraph contains the placeholder
                full_text = ''.join(run.text for run in paragraph.runs)
                if '{lead_pastor}' in full_text:
                    # This is the paragraph with our placeholder
                    # The text is split across 3 runs: "{", "lead_pastor", "}"
                    # We'll modify them in place to preserve all XML properties including glow effect
                    
                    if len(paragraph.runs) >= 3:
                        # Strategy: Modify the middle run to contain the full name,
                        # then clear the first and last runs
                        paragraph.runs[1].text = request.lead_pastor
                        paragraph.runs[0].text = ""
                        paragraph.runs[2].text = ""
                    elif len(paragraph.runs) == 1:
                        # If it's a single run, just replace the text
                        paragraph.runs[0].text = paragraph.runs[0].text.replace('{lead_pastor}', request.lead_pastor)
                    break
        
        # Save to a temporary file
        with tempfile.NamedTemporaryFile(
            suffix='.pptx',
            delete=False,
            prefix='message_for_all_generations_'
        ) as tmp_file:
            prs.save(tmp_file.name)
            tmp_path = tmp_file.name
        
        # Return the file
        return FileResponse(
            tmp_path,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            filename='message_for_all_generations_slide.pptx',
            headers={
                'Content-Disposition': 'attachment; filename="message_for_all_generations_slide.pptx"'
            }
        )
        
    except Exception as e:
        print(f"Error generating Message for All Generations slide: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Error generating Message for All Generations slide: {str(e)}"
        )