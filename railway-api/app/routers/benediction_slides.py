"""
Benediction slide generation endpoints
"""
import os
import tempfile
import base64
import io
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Pt, Inches
from PIL import Image
from app.core.schemas import BenedictionRequest
from pathlib import Path

router = APIRouter()


def create_benediction_slide(output_file, background_image=None, lead_pastor="Pastor"):
    """Create a benediction slide"""
    # Get the template path
    template_path = Path(__file__).parent.parent / "templates" / "benediction.pptx"
    if not template_path.exists():
        raise ValueError(f"Template not found: {template_path}")
    
    
    # Load the template presentation
    prs = Presentation(template_path)
    
    # Get the first slide
    if len(prs.slides) == 0:
        raise ValueError("Template has no slides")
    
    slide = prs.slides[0]
    
    # Find and update the {lead_pastor} placeholder - check all shapes
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = ''.join(run.text for run in paragraph.runs)
                if '{lead_pastor}' in full_text:
                    if len(paragraph.runs) >= 3:
                        paragraph.runs[1].text = lead_pastor
                        paragraph.runs[0].text = ""
                        paragraph.runs[2].text = ""
                    elif len(paragraph.runs) == 1:
                        paragraph.runs[0].text = paragraph.runs[0].text.replace('{lead_pastor}', lead_pastor)
                    break
    
    # Add the background image if provided
    if background_image:
        try:
            # Decode base64 image
            image_data = base64.b64decode(background_image.split(',')[1] if ',' in background_image else background_image)
            
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
                tmp_img.write(image_data)
                tmp_img_path = tmp_img.name
            
            # Get slide dimensions
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Add the background image covering the entire slide
            pic = slide.shapes.add_picture(tmp_img_path, 0, 0, slide_width, slide_height)
            
            # Move the picture to the back (z-order) - behind all existing shapes
            # First remove it from its current position (at the end)
            slide.shapes._spTree.remove(pic._element)
            # Then insert it at position 2 (after nvGrpSpPr at 0 and grpSpPr at 1)
            slide.shapes._spTree.insert(2, pic._element)
            
            # Clean up temp file
            os.unlink(tmp_img_path)
            
        except Exception as e:
            print(f"Error processing background image: {str(e)}")
    
    # Save the presentation
    prs.save(output_file)
    return output_file


@router.post("/api/generate-benediction-slide")
async def generate_benediction_slide(request: BenedictionRequest):
    """Generate a benediction slide with the lead pastor's name and background image"""
    try:
        # Get the template path
        template_path = Path(__file__).parent.parent / "templates" / "benediction.pptx"
        if not template_path.exists():
            raise HTTPException(status_code=500, detail=f"Template not found: {template_path}")
        
        
        # Load the template presentation
        prs = Presentation(template_path)
        
        # Get the first slide
        if len(prs.slides) == 0:
            raise HTTPException(
                status_code=500,
                detail="Template has no slides"
            )
        
        slide = prs.slides[0]
        
        # 1. First replace the {lead_pastor} text - check all shapes
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                # Find the paragraph with "{lead_pastor}" text
                for paragraph in shape.text_frame.paragraphs:
                    # Check if this paragraph contains the placeholder
                    full_text = ''.join(run.text for run in paragraph.runs)
                    if '{lead_pastor}' in full_text:
                        # The text is split across 3 runs: "{", "lead_pastor", "}"
                        # We'll modify them in place to preserve all XML properties including glow effect
                        
                        if len(paragraph.runs) >= 3:
                            # Strategy: Modify the middle run to contain the full name,
                            # then clear the first and last runs
                            paragraph.runs[1].text = request.lead_pastor
                            paragraph.runs[0].text = ""
                            paragraph.runs[2].text = ""
                        elif len(paragraph.runs) == 1:
                            # If it's a single run, just replace the text
                            paragraph.runs[0].text = paragraph.runs[0].text.replace('{lead_pastor}', request.lead_pastor)
                        break
        
        # 2. Add the background image as a new full-slide shape at the back
        if request.background_image:
            try:
                # Decode base64 image
                image_data = base64.b64decode(request.background_image.split(',')[1] if ',' in request.background_image else request.background_image)
                
                # Save to temporary file
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
                    tmp_img.write(image_data)
                    tmp_img_path = tmp_img.name
                
                # Get slide dimensions (standard 16:9 presentation)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # Add the background image covering the entire slide
                pic = slide.shapes.add_picture(tmp_img_path, 0, 0, slide_width, slide_height)
                
                # Move the picture to the back (z-order)
                # First remove it from its current position (at the end)
                slide.shapes._spTree.remove(pic._element)
                # Then insert it at position 2 (after nvGrpSpPr at 0 and grpSpPr at 1)
                slide.shapes._spTree.insert(2, pic._element)
                
                # Clean up temp file
                os.unlink(tmp_img_path)
                
            except Exception as e:
                print(f"Error processing background image: {str(e)}")
                # Continue without adding the background
        
        # Save to a temporary file
        with tempfile.NamedTemporaryFile(
            suffix='.pptx',
            delete=False,
            prefix='benediction_slide_'
        ) as tmp_file:
            prs.save(tmp_file.name)
            tmp_path = tmp_file.name
        
        # Return the file
        return FileResponse(
            tmp_path,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            filename='benediction_slide.pptx',
            headers={
                'Content-Disposition': 'attachment; filename="benediction_slide.pptx"'
            }
        )
        
    except Exception as e:
        print(f"Error generating benediction slide: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Error generating benediction slide: {str(e)}"
        )