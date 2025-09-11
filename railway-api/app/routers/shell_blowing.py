"""
Shell-blowing slides router for generating PowerPoint presentations
"""
import os
import shutil
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from typing import Dict, Any
from app.core.files import create_temp_file
from pathlib import Path
from pptx import Presentation

router = APIRouter()


@router.post("/generate-shell-blowing-slides")
async def generate_shell_blowing_slides(data: Dict[str, Any]):
    """Generate shell-blowing PowerPoint slides from template"""
    try:
        # Get the title from request (optional)
        title = data.get('title', 'Shell Blowing')
        
        # Path to the template file
        template_path = Path(__file__).parent.parent / "templates" / "shell-blowing.pptx"
        if not template_path.exists():
            raise HTTPException(status_code=500, detail=f"Template not found: {template_path}")
        
        # Create temporary output file
        output_path = create_temp_file(suffix='.pptx')
        
        # Copy the template to the output file
        shutil.copy2(template_path, output_path)
        
        # If we want to modify the title, we can open and edit the presentation
        if title != 'Shell Blowing':
            prs = Presentation(output_path)
            
            # Modify the first slide's title if it exists
            if prs.slides:
                slide = prs.slides[0]
                # Look for title placeholder
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        # Replace "Shell Blowing" text with custom title
                        if 'Shell Blowing' in shape.text_frame.text:
                            shape.text_frame.text = shape.text_frame.text.replace(
                                'Shell Blowing', title
                            )
                
                # Save the modified presentation
                prs.save(output_path)
        
        # Return the file
        return FileResponse(
            path=output_path,
            filename=f"shell_blowing_{title.replace(' ', '_').lower()}.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))